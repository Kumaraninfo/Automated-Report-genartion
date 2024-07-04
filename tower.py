from pptx import Presentation
from pptx.util import Inches, Pt
import json
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json


 
pr1 = Presentation('templates.pptx')


f = open('new.json')
data = json.load(f)
n=0
a=0
c=0
slide_number = 0


#TABLE
for i in data['data']: #data

    slide = pr1.slides[0]  
    
    left_table = Inches(1)
    top_table = Inches(3)
    width_table = Inches(6)
    height_table = Inches(3) 
    table1_frame = slide.shapes.add_table(10,3, left_table,top_table,width_table,height_table)
    table1 = table1_frame.table
   
    slide_number += 1

    # Add slide number
    slide_number_box = slide.shapes.add_textbox(Inches(5.85), Inches(10.17), Inches(1), Inches(0.5))
    text_frame = slide_number_box.text_frame
    text_frame.text = f"Page No : {slide_number}"
    font = text_frame.paragraphs[0].runs[0].font
    
    font.bold = True
    font.size = Pt(16)


    for j in data['data'][i]['towerJsonFilePath']: #data>03gLs1



        #Heading Of The Table
        left = Inches(2.3)
        top = Inches(2.1)
        width = Inches(6)
        height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Transmission Compound Index"
        text_box.width = Inches(7)
        text_box.height = Inches(1.5)
        font = text_frame.paragraphs[0].runs[0].font
        font.name = "Calibri"
        font.font_family = "Body"
        font.color.rgb = RGBColor(12, 53, 106)  
        font.bold = True


        cell = table1.cell(0, 0)
        cell.text = "Components Inspected"
        cell_font = cell.text_frame.paragraphs[0].runs[0].font
        cell_font.name = "Calibri (Body)"
        cell_font.font_family = "Body"
        cell_font.size = Pt(14)
        cell_font.bold = False  # Set bold to False
        cell_font.color.rgb = RGBColor(0, 0, 0)  # Black text color
        cell_text_frame = cell.text_frame
        cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)

        cell = table1.cell(0,1)
        cell.text = "Total Components Count"
        cell_font = cell.text_frame.paragraphs[0].runs[0].font
        cell_font.name = "Calibri (Body)"
        cell_font.font_family = "Body"
        cell_font.size = Pt(14)
        cell_font.bold = False
        cell_font.color.rgb = RGBColor(0, 0, 0)  # Black text color
        cell_text_frame = cell.text_frame
        cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)

        cell = table1.cell(0,2)
        cell.text = "Components Defect Count"
        cell_font = cell.text_frame.paragraphs[0].runs[0].font
        cell_font.name = "Calibri (Body)"
        cell_font.font_family = "Body"
        cell_font.size = Pt(14)
        cell_font.bold = False
        cell_font.color.rgb = RGBColor(0, 0, 0)  # Black text color
        cell_text_frame = cell.text_frame
        cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)

        for z in data['data'][i]['towerJsonFilePath'][j]: #data>03gLs1>towerJsonFilePath
            # subcomponent name 
            c=c+1
            for w in data['data'][i]['towerJsonFilePath'][j][z]:
                a=a+1
               
                if data['data'][i]['towerJsonFilePath'][j][z][w]['isDefect']==True:
                    n=n+1
            # print(z)    
            cell = table1.cell(c, 0)
            cell.text = z.capitalize()
            cell_font = cell.text_frame.paragraphs[0].runs[0].font
            cell_font.name = "Calibri (Body)"
            cell_font.font_family = "Body"
            cell_font.size = Pt(14)
            cell_text_frame = cell.text_frame
            cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(144, 183, 222)

            cell = table1.cell(c, 1)
            cell.text = str(a)
            cell_font = cell.text_frame.paragraphs[0].runs[0].font
            cell_font.name = "Calibri (Body)"
            cell_font.font_family = "Body"
            cell_font.size = Pt(14)
            cell_text_frame = cell.text_frame
            cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(144, 183, 222)

            cell = table1.cell(c, 2)
            cell.text = str(n)
            cell_font = cell.text_frame.paragraphs[0].runs[0].font
            cell_font.name = "Calibri (Body)"
            cell_font.font_family = "Body"
            cell_font.size = Pt(14)
            cell_text_frame = cell.text_frame
            cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(144, 183, 222)
            
            
            
            a=0
            n=0
        
        
        c=0
            # break #control foue
        break
    break
                


components_others =[]
components_ins_jum =[]


for mainkey, mainvalue in data.items():
    for key_1, value_1 in mainvalue.items():
        for key_2, value_2 in value_1['towerJsonFilePath'].items():
            for key_3, value_3 in value_2.items():
                if key_3 in ['anticlimb', 'arcinghorn', 'foundationleg' , 'dangersign','vibrationdamper','gwjumper','crossarm']:
                    for key_4, value_4 in value_3.items():

                        components_others.append(value_4['Name'])

                        if len(components_others) == 1:
                            slide_number += 1
                            slide = pr1.slides.add_slide(pr1.slide_layouts[0]) 


                            # Add slide number
                            slide_number_box = slide.shapes.add_textbox(Inches(5.85), Inches(10.17), Inches(1), Inches(0.5))
                            text_frame = slide_number_box.text_frame
                            text_frame.text = f"Page No : {slide_number}"
                            font = text_frame.paragraphs[0].runs[0].font
                            font.bold = True
                            
                            

                            top_position = Inches(1.5)
                                
                    

                               
                            for i, component_name in enumerate(components_others):
                                left = Inches(2.6)
                                width = Inches(7)
                                height = Inches(1.0)
                                text_box = slide.shapes.add_textbox(left, top_position, width, height)
                                text_frame = text_box.text_frame
                                text_frame.text = component_name
                                font = text_frame.paragraphs[0].runs[0].font
                                font.name = "Calibri"
                                font.font_family = "Body"
                                font.color.rgb = RGBColor(12, 53, 106)
                                font.bold = True
                                font.size = Pt(22)

                                image_path_center = "pic2.png"
                                left_image_right = Inches(2.2) 
                                top_image_right = top_position + Inches(1.0)  
                                width_image_right = Inches(3.334645669)
                                height_image_right = Inches(2.5)
                                slide.shapes.add_picture(image_path_center, left_image_right, top_image_right, width_image_right, height_image_right)

                                top_position += height + Inches(0.2)  


                            zoom_value = len(value_4['Zoomimage'])
                            if value_4['Resolved']['Isresolved'] == True:
                               
                            
                               
                                
                                # Add slides and component information
                                if zoom_value == 1:
                                    

                                    
                                    
                                    
                                    # Add image
                                    image_path_center = value_4['Zoomimage'][0]
                                    left_image_center = Inches(2.5)
                                    top_image_center = Inches(5.5)
                                    width_image_center = Inches(2.5)
                                    height_image_center = Inches(2.5)
                                    slide.shapes.add_picture(image_path_center, left_image_center, top_image_center, width_image_center, height_image_center)
                                    
                                    
                                    left = Inches(1.5)
                                    top = Inches(9.2)
                                    width = Inches(6)
                                    height = Inches(1)
                                    text_box = slide.shapes.add_textbox(left, top, width, height)
                                    text_frame = text_box.text_frame
                                    text_frame.text = value_4['CommentType']
                                    text_box.width = Inches(7)
                                    text_box.height = Inches(1.5)
                                    font = text_frame.paragraphs[0].runs[0].font
                                    font.name = "Calibri"
                                    font.font_family = "Body"
                                    font.color.rgb = RGBColor(12, 53, 106)  
                                    font.bold = True

                            elif zoom_value > 1:
                                slide = pr1.slides.add_slide(pr1.slide_layouts[0])
                                slide_number += 1
                                
                                
                                

                                left = Inches(1.5)
                                top = Inches(9)
                                width = Inches(6)
                                height = Inches(1)
                                text_box = slide.shapes.add_textbox(left, top, width, height)
                                text_frame = text_box.text_frame
                                text_frame.text = value_4['CommentType']
                                text_box.width = Inches(7)
                                text_box.height = Inches(1.5)
                                font = text_frame.paragraphs[0].runs[0].font
                                font.name = "Calibri"
                                font.font_family = "Body"
                                font.color.rgb = RGBColor(12, 53, 106)  
                                font.bold = True
                                
                                # Add images in a grid layout
                                left_start = Inches(1)
                                top_start = Inches(4)
                                image_width = Inches(2.5)
                                image_height = Inches(2.5)
                                spacing = Inches(0.35)

                                for idx, image_path in enumerate(value_4['Zoomimage']):
                                    col = idx % 2  # 2 images per row
                                    row = idx // 2  # Rows based on index
                                    left = left_start + col * (image_width + spacing)
                                    top = top_start + row * (image_height + spacing)
                                    slide.shapes.add_picture(image_path, left, top, image_width, image_height)
                        
                        components_others = [] 

                if key_3 in ['insulator', 'jumper']:
                    for key_4, value_4 in value_3.items():
                                   
                        components_ins_jum.append(value_4)
                                    
                        if len(components_ins_jum) == 1:
                            slide_number += 1
                            slide = pr1.slides.add_slide(pr1.slide_layouts[0])

                            # Add slide number
                            slide_number_box = slide.shapes.add_textbox(Inches(5.85), Inches(10.17), Inches(1), Inches(0.5))
                            text_frame = slide_number_box.text_frame
                            text_frame.text = f"Page No : {slide_number}"
                            font = text_frame.paragraphs[0].runs[0].font
                            font.bold = True
                            font.size = Pt(16)
                            
                            
                                

                                
                                
                                

                                
                            top_position = Inches(1.1)

                                
                            for component_data in components_ins_jum:
                                
                                left = Inches(2.7)
                                width = Inches(7)
                                height = Inches(1.0)
                                text_box = slide.shapes.add_textbox(left, top_position, width, height)
                                text_frame = text_box.text_frame
                                text_frame.text = component_data['Name']
                                font = text_frame.paragraphs[0].runs[0].font
                                font.name = "Calibri"
                                font.font_family = "Body"
                                font.color.rgb = RGBColor(12, 53, 106)
                                font.bold = True
                                font.size = Pt(22)

                                # Determine which fields to display based on component type
                                if key_3 == 'jumper':
                                    fields = ['sp1', 'sp2', 'sp3']
                                    image_path_left = "Normal.png"  
                                    image_path_right = "thermal.png"  
                                elif key_3 == 'insulator':
                                    fields = [  'min','avg','max']
                                    image_path_left = "normal.png"  
                                    image_path_right = "thermal.png"  
                                else:
                                    fields = []
                                    image_path_left = None
                                    image_path_right = None

                                # Add image at the top-left of the table
                                if image_path_left:
                                    left_image_left = Inches(0.6)
                                    top_image_left = top_position + Inches(0.7)  
                                    width_image_left = Inches(3)
                                    height_image_left = Inches(2.5)
                                    slide.shapes.add_picture(image_path_left, left_image_left, top_image_left, width_image_left, height_image_left)

                                # Add image at the top-right of the table
                                if image_path_right:
                                    left_image_right = Inches(3.85)  
                                    top_image_right = top_position + Inches(0.7)  
                                    width_image_right = Inches(3)
                                    height_image_right = Inches(2.5)
                                    slide.shapes.add_picture(image_path_right, left_image_right, top_image_right, width_image_right, height_image_right)

                                # Add table for the component
                                left_table = Inches(0.7)
                                top_table = top_position + Inches(3.5)
                                width_table = Inches(6)
                                height_table = Inches(1)
                                table_frame = slide.shapes.add_table(2, 5, left_table, top_table, width_table, height_table).table

                                # Set table headings based on component type
                                if key_3 == 'jumper':
                                    table_frame.cell(0, 0).text = "Compound"
                                    table_frame.cell(0, 1).text = "Sp1"
                                    table_frame.cell(0, 2).text = "Sp2"
                                    table_frame.cell(0, 3).text = "Sp3"
                                    table_frame.cell(0, 4).text = "Status"
                                elif key_3 == 'insulator':
                                    table_frame.cell(0, 0).text = "Compound"
                                    table_frame.cell(0, 1).text = "Min Temp"
                                    table_frame.cell(0, 2).text = "Aug Temp"
                                    table_frame.cell(0, 3).text = "Max Temp"
                                    table_frame.cell(0, 4).text = "Status"

                                table_frame.cell(1, 0).text = component_data['Name']
                                # Ensure to check if the field exists in the data before accessing it
                                table_frame.cell(1, 1).text = component_data['TempValues'].get(fields[0], "") if fields and len(fields) > 0 else ""
                                table_frame.cell(1, 2).text = component_data['TempValues'].get(fields[1], "") if fields and len(fields) > 1 else ""
                                table_frame.cell(1, 3).text = component_data['TempValues'].get(fields[2], "") if fields and len(fields) > 2 else ""
                                table_frame.cell(1, 4).text = component_data['WorkStatusAndQC']['Status']

                                # Format table cells
                                for row in range(2):
                                    for col in range(5):
                                        cell = table_frame.cell(row, col)
                                        if len(cell.text_frame.paragraphs) > 0 and len(cell.text_frame.paragraphs[0].runs) > 0:
                                            cell_font = cell.text_frame.paragraphs[0].runs[0].font
                                            cell_font.name = "Calibri (Body)"
                                            cell_font.font_family = "Body"
                                            cell_font.size = Pt(14)
                                            cell_text_frame = cell.text_frame
                                            cell_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                
                                zoom_value = len(value_4['Zoomimage'])
                                if value_4['Resolved']['Isresolved'] == True:
                                   
                                # Add slides and component information
                                    if zoom_value == 1:
                                        image_path_center = value_4['Zoomimage'][0]
                                        left_image_center = Inches(2.5)
                                        top_image_center = Inches(5.9)
                                        width_image_center = Inches(2.5)
                                        height_image_center = Inches(2.5)
                                        slide.shapes.add_picture(image_path_center, left_image_center, top_image_center, width_image_center, height_image_center)
                                        
                                        
                                        left = Inches(1.5)
                                        top = Inches(9)
                                        width = Inches(6)
                                        height = Inches(1)
                                        text_box = slide.shapes.add_textbox(left, top, width, height)
                                        text_frame = text_box.text_frame
                                        text_frame.text = value_4['CommentType']
                                        text_box.width = Inches(7)
                                        text_box.height = Inches(1.5)
                                        font = text_frame.paragraphs[0].runs[0].font
                                        font.name = "Calibri"
                                        font.font_family = "Body"
                                        font.color.rgb = RGBColor(12, 53, 106)  
                                        font.bold = True

                                    elif zoom_value > 1:
                                        
                                        slide_number += 1
                                    
                                       

                                        left = Inches(1.5)
                                        top = Inches(9)
                                        width = Inches(6)
                                        height = Inches(1)
                                        text_box = slide.shapes.add_textbox(left, top, width, height)
                                        text_frame = text_box.text_frame
                                        text_frame.text = value_4['CommentType']
                                        text_box.width = Inches(7)
                                        text_box.height = Inches(1.5)
                                        font = text_frame.paragraphs[0].runs[0].font
                                        font.name = "Calibri"
                                        font.font_family = "Body"
                                        font.color.rgb = RGBColor(12, 53, 106)  
                                        font.bold = True
                                        
                                        # Add images in a grid layout
                                        left_start = Inches(1)
                                        top_start = Inches(6)
                                        image_width = Inches(2.5)
                                        image_height = Inches(2.5)
                                        spacing = Inches(0.35)

                                        for idx, image_path in enumerate(value_4['Zoomimage']):
                                            col = idx % 2  # 2 images per row
                                            row = idx // 2  # Rows based on index
                                            left = left_start + col * (image_width + spacing)
                                            top = top_start + row * (image_height + spacing)
                                            slide.shapes.add_picture(image_path, left, top, image_width, image_height)
                            
                            components_ins_jum = []
        break                
                         
                    

pr1.save("newdesign.pptx")